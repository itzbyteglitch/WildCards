#!/usr/bin/env python3
"""Generate a professional PowerPoint presentation for the WildCards project.

WildCards — a production-quality, browser-based multiplayer UNO game.
Tech: React 19 · TypeScript · TanStack Start · Cloudflare Workers · D1 · WebSockets.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

# ----------------------------------------------------------------------------
# Theme constants
# ----------------------------------------------------------------------------
BG          = RGBColor(0x0B, 0x12, 0x20)   # deep slate background
BG_SOFT     = RGBColor(0x12, 0x1B, 0x2E)   # panel background
BG_CARD     = RGBColor(0x1A, 0x25, 0x3A)   # card background
BORDER      = RGBColor(0x2A, 0x38, 0x54)   # subtle border
TEXT        = RGBColor(0xF1, 0xF5, 0xF9)   # near-white text
MUTED       = RGBColor(0x94, 0xA3, 0xB8)   # slate-400
FAINT       = RGBColor(0x64, 0x74, 0x8B)   # slate-500

RED         = RGBColor(0xE5, 0x39, 0x35)   # UNO red
YELLOW      = RGBColor(0xFD, 0xD8, 0x35)   # UNO yellow
GREEN       = RGBColor(0x43, 0xA0, 0x47)   # UNO green
BLUE        = RGBColor(0x1E, 0x88, 0xE5)   # UNO blue
ACCENT      = RED

FONT        = "Segoe UI"
FONT_LIGHT  = "Segoe UI Light"
FONT_MONO   = "Consolas"

SLIDE_W     = Inches(13.333)
SLIDE_H     = Inches(7.5)

ASSETS = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "docs", "presentation", "assets")
OUT    = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "docs", "presentation", "Project_Presentation.pptx")

os.makedirs(ASSETS, exist_ok=True)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]

PAGE = [0]

# ----------------------------------------------------------------------------
# Helpers
# ----------------------------------------------------------------------------

def new_slide():
    PAGE[0] += 1
    return prs.slides.add_slide(BLANK)


def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=0.75, shape=MSO_SHAPE.RECTANGLE, radius=None, shadow=True):
    sp = slide.shapes.add_shape(shape, x, y, w, h)
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    return sp


def add_text(slide, x, y, w, h, runs, size=16, color=TEXT, bold=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT,
             line_spacing=1.0, space_after=0, space_before=0):
    """runs: str OR list of paragraphs; each paragraph is str OR list of (text, dict) run tuples."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0

    if isinstance(runs, str):
        runs = [runs]
    first = True
    for para in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        p.space_before = Pt(space_before)
        if isinstance(para, str):
            para = [(para, {})]
        for txt, style in para:
            r = p.add_run()
            r.text = txt
            r.font.name = style.get("font", font)
            r.font.size = Pt(style.get("size", size))
            r.font.bold = style.get("bold", bold)
            r.font.italic = style.get("italic", False)
            r.font.color.rgb = style.get("color", color)
    return tb


def color_stripe(slide):
    """Four-color UNO stripe across the very top."""
    q = int(SLIDE_W) // 4
    for i, c in enumerate([RED, YELLOW, GREEN, BLUE]):
        add_rect(slide, Emu(q * i), 0, Emu(q), Inches(0.09), fill=c)


def header(slide, title, kicker=None, accent=ACCENT):
    color_stripe(slide)
    if kicker:
        add_text(slide, Inches(0.65), Inches(0.42), Inches(9.5), Inches(0.3),
                 kicker.upper(), size=11, color=accent, bold=True)
        ty = Inches(0.72)
    else:
        ty = Inches(0.5)
    add_text(slide, Inches(0.65), ty, Inches(11.8), Inches(0.75),
             title, size=30, color=TEXT, bold=True, font=FONT)
    add_rect(slide, Inches(0.67), ty + Inches(0.62), Inches(1.35), Inches(0.045), fill=accent)
    # Footer
    add_text(slide, Inches(0.65), Inches(7.12), Inches(6.5), Inches(0.3),
             "WILD CARDS  ·  Technical Project Presentation", size=9, color=FAINT)
    add_text(slide, Inches(12.35), Inches(7.12), Inches(0.6), Inches(0.3),
             str(PAGE[0]), size=9, color=FAINT, align=PP_ALIGN.RIGHT)


def arrow(slide, x1, y1, x2, y2, color=FAINT, w=2.25):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    conn.line.color.rgb = color
    conn.line.width = Pt(w)
    conn.shadow.inherit = False
    ln = conn._element.spPr.find(qn('a:ln'))
    if ln is not None:
        tail = ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'})
        ln.append(tail)
    return conn


def chip(slide, x, y, w, text, color, filled=True, size=11):
    sp = add_rect(slide, x, y, w, Inches(0.34),
                  fill=color if filled else BG_SOFT,
                  line=None if filled else color,
                  shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
    tf = sp.text_frame
    tf.word_wrap = False
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = FONT_MONO if filled else FONT
    r.font.size = Pt(size)
    r.font.bold = not filled
    r.font.color.rgb = BG if filled else color
    return sp


def card(slide, x, y, w, h, title, body, accent=ACCENT, title_size=15, body_size=12):
    add_rect(slide, x, y, w, h, fill=BG_CARD, line=BORDER, line_w=1.0,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.055)
    add_rect(slide, x, y, Inches(0.07), h, fill=accent)
    add_text(slide, x + Inches(0.28), y + Inches(0.18), w - Inches(0.5), Inches(0.4),
             title, size=title_size, color=TEXT, bold=True)
    add_text(slide, x + Inches(0.28), y + Inches(0.62), w - Inches(0.5), h - Inches(0.8),
             body, size=body_size, color=MUTED, line_spacing=1.15)


def step_badge(slide, x, y, num, color=ACCENT):
    add_rect(slide, x, y, Inches(0.34), Inches(0.34), fill=color, shape=MSO_SHAPE.OVAL)
    add_text(slide, x, y + Inches(0.015), Inches(0.34), Inches(0.3),
             str(num), size=13, color=BG, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def image_centered(slide, path, x, y, w):
    from PIL import Image
    im = Image.open(path)
    ratio = im.size[0] / im.size[1]
    h = w / ratio
    return slide.shapes.add_picture(path, x, y, w, h), h


def bullets(slide, x, y, w, h, items, size=14, gap=8, lead_color=TEXT, body_color=MUTED):
    paras = []
    for it in items:
        if isinstance(it, tuple):
            lead, rest = it
            paras.append([(lead, {"bold": True, "color": lead_color}), (rest, {"color": body_color})])
        else:
            paras.append([("▸  ", {"bold": True, "color": ACCENT}), (it, {"color": body_color})])
    add_text(slide, x, y, w, h, paras, size=size, line_spacing=1.12, space_after=gap)


# ============================================================================
# SLIDE 1 — TITLE
# ============================================================================
s = new_slide()
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=BG)
# soft radial glow behind cards
add_rect(s, Inches(7.2), Inches(0.4), Inches(6.2), Inches(5.6), fill=BG_SOFT, shape=MSO_SHAPE.OVAL)

# Fan of UNO cards (top-right)
fan = [
    (RED, -22, Inches(8.7), Inches(1.15)),
    (YELLOW, -8, Inches(9.1), Inches(1.15)),
    (GREEN, 8, Inches(9.1), Inches(1.15)),
    (BLUE, 22, Inches(8.7), Inches(1.15)),
]
for i, (c, rot, x, y) in enumerate(fan):
    cw = Inches(1.15)
    ch = Inches(1.75)
    sp = add_rect(s, x, y, cw, ch, fill=c, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.14)
    sp.rotation = rot
    # inner card
    inner = add_rect(s, x + Inches(0.14), y + Inches(0.22), cw - Inches(0.28), ch - Inches(0.44),
                     fill=BG, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.18)
    inner.rotation = rot
    # value text on inner
    add_text(s, x + Inches(0.14), y + Inches(0.5), cw - Inches(0.28), Inches(0.6),
             "UNO", size=22, color=c, bold=True, align=PP_ALIGN.CENTER)

# Title block (left)
add_text(s, Inches(0.9), Inches(1.6), Inches(7.2), Inches(0.4),
         "TECHNICAL PROJECT PRESENTATION", size=13, color=YELLOW, bold=True)

add_text(s, Inches(0.88), Inches(2.0), Inches(8.0), Inches(1.5),
         [[("WILD ", {"size": 66, "bold": True, "color": TEXT}),
           ("CARDS", {"size": 66, "bold": True, "color": RED})]],
         font=FONT)

add_text(s, Inches(0.9), Inches(3.15), Inches(8.0), Inches(0.6),
         "Browser Multiplayer UNO — Reimagined for the Web",
         size=22, color=BLUE, bold=True)

add_text(s, Inches(0.9), Inches(3.9), Inches(7.6), Inches(1.2),
         "A production-quality, real-time card game playable instantly in any "
         "browser — no downloads, no accounts. Powered by a server-authoritative "
         "engine and a globally distributed realtime layer.",
         size=14, color=MUTED, line_spacing=1.25)

# Tech chips
chips = ["React 19", "TypeScript", "TanStack Start", "Cloudflare Workers", "Durable Objects", "WebSockets", "D1 SQLite"]
cx = Inches(0.9)
cy = Inches(5.35)
row_h = Inches(0.34)
for i, t in enumerate(chips):
    w = Inches(0.36 + 0.082 * len(t))
    if cx + w > Inches(8.6):
        cx = Inches(0.9)
        cy += row_h + Inches(0.14)
    chip(s, cx, cy, w, t, BLUE if i % 2 == 0 else GREEN)
    cx += w + Inches(0.18)

add_text(s, Inches(0.9), Inches(6.55), Inches(8.0), Inches(0.4),
         "2–8 players  ·  private rooms  ·  smart bots  ·  cross-device realtime",
         size=12, color=FAINT)

# ============================================================================
# SLIDE 2 — AGENDA
# ============================================================================
s = new_slide()
header(s, "Agenda", kicker="Presentation outline")
agenda = [
    ("01  Overview", "What WildCards is and why it exists"),
    ("02  Problem & Objectives", "The gaps it addresses and the goals set"),
    ("03  Key Features", "What players actually get"),
    ("04  Technology Stack", "Every layer, framework and library"),
    ("05  System Architecture", "How the pieces fit together"),
    ("06  Realtime Transport", "The auto-selecting delivery layer"),
    ("07  Game Engine", "Server-authoritative UNO rules"),
    ("08  Data, Auth & Multiplayer Flow", "Sessions, persistence, live play"),
    ("09  Deployment & CI/CD", "Shipping to the edge"),
    ("10  Performance, Security & Challenges", "Hard problems, real answers"),
    ("11  Demo & Next Steps", "See it live, then what's next"),
]
for i, (t, d) in enumerate(agenda):
    col = i // 6
    row = i % 6
    x = Inches(0.9) + col * Inches(6.2)
    y = Inches(1.55) + row * Inches(0.86)
    step_badge(s, x, y + Inches(0.02), i + 1, color=[RED, YELLOW, GREEN, BLUE][i % 4])
    add_text(s, x + Inches(0.5), y, Inches(5.4), Inches(0.35), t, size=16, color=TEXT, bold=True)
    add_text(s, x + Inches(0.5), y + Inches(0.34), Inches(5.4), Inches(0.35), d, size=11.5, color=MUTED)

# ============================================================================
# SLIDE 3 — PROJECT OVERVIEW
# ============================================================================
s = new_slide()
header(s, "Project Overview", kicker="What is WildCards")
add_text(s, Inches(0.9), Inches(1.6), Inches(11.5), Inches(1.5),
         [[("WildCards ", {"bold": True, "color": TEXT}),
           ("is a production-quality, browser-based multiplayer implementation of the classic card game UNO. "
            "A single invite link drops players straight into a live game — the room is created instantly, "
            "bots fill empty seats, and every move is validated by an authoritative server so the rules can "
            "never be cheated.", {"color": MUTED})]],
         size=15, line_spacing=1.3)

stats = [
    ("2–8", "players per room", RED),
    ("≈30s", "turn timers keep games moving", YELLOW),
    ("100%", "server-authoritative rules", GREEN),
    ("0", "accounts or downloads required", BLUE),
]
for i, (num, label, c) in enumerate(stats):
    x = Inches(0.9) + i * Inches(3.0)
    card(s, x, Inches(3.35), Inches(2.7), Inches(1.7), "", "", accent=c)
    add_text(s, x + Inches(0.28), Inches(3.55), Inches(2.2), Inches(0.7),
             num, size=30, color=c, bold=True)
    add_text(s, x + Inches(0.28), Inches(4.3), Inches(2.25), Inches(0.6),
             label, size=11.5, color=MUTED, line_spacing=1.1)

card(s, Inches(0.9), Inches(5.45), Inches(11.5), Inches(1.25), "One codebase, three environments",
     "The frontend runs as a static site; TanStack Start powers SSR and type-safe server functions; "
     "a Cloudflare Worker + Durable Object hosts each room as a globally addressable relay. The same "
     "code deploys across tabs, browsers and devices.", accent=YELLOW)

# ============================================================================
# SLIDE 4 — PROBLEM STATEMENT
# ============================================================================
s = new_slide()
header(s, "Problem Statement", kicker="Why rebuild a classic game")
probs = [
    ("Physical UNO is slow", "Setup, shuffling, rule disputes and scorekeeping interrupt the fun. "
     "Games can't happen without everyone in the same room."),
    ("Real-time multiplayer is hard", "State sync, latency, reconnects and fairness are notoriously "
     "difficult to get right — most web games cheat by trusting the client."),
    ("Most web clones feel broken", "Dead lobbies, logins required, ads, slow servers, or single-device "
     "demos that fall apart under real concurrent play."),
    ("Trusting the client breaks the game", "A client-authoritative card game can be modified to draw "
     "cards, see hands, or win instantly — fair play demands the server make every decision."),
]
for i, (t, d) in enumerate(probs):
    x = Inches(0.9) + (i % 2) * Inches(6.0)
    y = Inches(1.65) + (i // 2) * Inches(2.55)
    step_badge(s, x, y, i + 1, color=[RED, YELLOW, GREEN, BLUE][i])
    add_text(s, x + Inches(0.55), y - Inches(0.02), Inches(5.3), Inches(0.4), t, size=17, color=TEXT, bold=True)
    add_text(s, x + Inches(0.55), y + Inches(0.42), Inches(5.2), Inches(1.9), d, size=12.5, color=MUTED, line_spacing=1.25)

# ============================================================================
# SLIDE 5 — OBJECTIVES
# ============================================================================
s = new_slide()
header(s, "Project Objectives", kicker="Goals that shaped every decision")
objs = [
    ("Instant, frictionless play", "One click to start against bots; one invite link to gather friends. "
     "No signup, no downloads, no installs."),
    ("Server-authoritative fairness", "Cards are dealt, shuffled and validated on the server. Clients "
     "only ever receive their own hand."),
    ("Rock-solid realtime", "WebSockets by default, graceful fallbacks, reconnection, and a live poller "
     "that keeps every client converged on the same state."),
    ("A complete UNO ruleset", "Draw-2 / Wild+4 stacking chains, missed-UNO penalties, skips, reverses, "
     "wild color selection, and automatic turn timeouts."),
    ("Professional engineering", "Strict TypeScript, deterministic engine, CI pipeline with typecheck + "
     "build + lint, and a deployable edge architecture."),
]
for i, (t, d) in enumerate(objs):
    y = Inches(1.55) + i * Inches(1.06)
    card(s, Inches(0.9), y, Inches(11.5), Inches(0.92), "", "", accent=[RED, YELLOW, GREEN, BLUE][i % 4])
    add_text(s, Inches(1.25), y + Inches(0.12), Inches(3.6), Inches(0.7),
             [[("0%d  " % (i + 1), {"bold": True, "color": [RED, YELLOW, GREEN, BLUE][i % 4]}),
               (t, {"bold": True, "color": TEXT})]],
             size=14.5, line_spacing=1.1)
    add_text(s, Inches(4.95), y + Inches(0.12), Inches(7.2), Inches(0.75), d, size=12, color=MUTED, line_spacing=1.12)

# ============================================================================
# SLIDE 6 — KEY FEATURES
# ============================================================================
s = new_slide()
header(s, "Key Features", kicker="What players get")
feats = [
    ("Instant play vs smart bots", "Start a match in one click; intelligent bots fill empty seats and "
     "follow the same rules as humans.", RED),
    ("Private rooms & invite links", "A 6-letter code and a shareable URL — up to 8 players join from "
     "any device or network.", YELLOW),
    ("Cross-device realtime", "Auto-selecting transport keeps every player in sync with live turn "
     "timers and a move log.", GREEN),
    ("Full UNO rule engine", "Stacking +2 / +4 chains, missed-UNO penalties, skips, reverses, wild "
     "color picker, turn timeouts.", BLUE),
    ("Stats & leaderboard", "Guest profiles track wins, score and history with a global leaderboard — "
     "no account required.", RED),
    ("Polished, animated UI", "Framer-Motion card physics, glassmorphism, dark theme, and full "
     "responsive layouts.", YELLOW),
]
for i, (t, d, c) in enumerate(feats):
    x = Inches(0.9) + (i % 3) * Inches(3.95)
    y = Inches(1.6) + (i // 3) * Inches(2.7)
    card(s, x, y, Inches(3.65), Inches(2.4), t, d, accent=c, title_size=14, body_size=11.5)

# ============================================================================
# SLIDE 7 — TECHNOLOGY STACK
# ============================================================================
s = new_slide()
header(s, "Technology Stack", kicker="Every layer, framework and library")

# Frontend column
add_text(s, Inches(0.9), Inches(1.55), Inches(6), Inches(0.4), "Frontend", size=18, color=RED, bold=True)
front = ["React 19", "TypeScript 5", "Vite 8", "TanStack Router", "TanStack Start", "TanStack Query",
         "Zustand", "Tailwind CSS 4", "Framer Motion", "Radix UI", "Zod", "Recharts"]
cx, cy = Inches(0.9), Inches(2.05)
for t in front:
    w = Inches(0.5 + 0.078 * len(t))
    if cx + w > Inches(6.9):
        cx = Inches(0.9)
        cy += Inches(0.46)
    chip(s, cx, cy, w, t, BLUE, filled=False)
    cx += w + Inches(0.18)

# Backend column
add_text(s, Inches(7.4), Inches(1.55), Inches(6), Inches(0.4), "Backend & Infrastructure", size=18, color=GREEN, bold=True)
back = ["Cloudflare Workers", "Durable Objects", "D1 (SQLite)", "Hono", "WebSockets", "Supabase Realtime",
        "BroadcastChannel", "JWT auth", "Wrangler CLI", "GitHub Actions"]
cx, cy = Inches(7.4), Inches(2.05)
for t in back:
    w = Inches(0.5 + 0.078 * len(t))
    if cx + w > Inches(12.4):
        cx = Inches(7.4)
        cy += Inches(0.46)
    chip(s, cx, cy, w, t, GREEN, filled=False)
    cx += w + Inches(0.18)

# Highlight card
card(s, Inches(0.9), Inches(5.3), Inches(11.5), Inches(1.45), "Game engine: pure, deterministic TypeScript",
     "The core UNO logic (deck building, shuffling, turn order, stacking, scoring) lives in a zero-dependency "
     "TypeScript module. Being deterministic, it can run identically on the server for validation and in "
     "tests for verification.", accent=YELLOW)

# ============================================================================
# SLIDE 8 — SYSTEM ARCHITECTURE
# ============================================================================
s = new_slide()
header(s, "System Architecture", kicker="How the pieces fit together")

def arch_box(x, y, w, h, title, sub, c, fill=BG_CARD):
    add_rect(s, x, y, w, h, fill=fill, line=c, line_w=1.5, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.12)
    add_text(s, x, y + Inches(0.14), w, Inches(0.4), title, size=14, color=c, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.12), y + Inches(0.52), w - Inches(0.24), h - Inches(0.6),
             sub, size=9.5, color=MUTED, align=PP_ALIGN.CENTER, line_spacing=1.1)

# Client
arch_box(Inches(0.9), Inches(1.6), Inches(3.0), Inches(1.45), "CLIENT", "Browser app\nReact 19 UI · TanStack Router\nreads server-authoritative state", BLUE)
# Server
arch_box(Inches(5.2), Inches(1.6), Inches(3.3), Inches(1.45), "TANSTACK START SERVER", "SSR + type-safe Server Functions\nREST actions · guest auth\nwrites to room & DB", GREEN)
# Room
arch_box(Inches(9.2), Inches(1.6), Inches(3.2), Inches(1.45), "ROOM SERVER", "Cloudflare Worker + Durable Object\nWebSocket relay · roster\nauthoritative state snapshot", RED)

# DB + Auth row
arch_box(Inches(5.2), Inches(4.1), Inches(3.3), Inches(1.25), "D1 · SQLITE", "profiles · game history\nleaderboard · stats", YELLOW)
arch_box(Inches(9.2), Inches(4.1), Inches(3.2), Inches(1.25), "AUTH", "guest sessions · JWT\nlocalStorage persistence", YELLOW)

# arrows
arrow(s, Inches(3.9), Inches(2.3), Inches(5.2), Inches(2.3), color=FAINT)
arrow(s, Inches(8.5), Inches(2.3), Inches(9.2), Inches(2.3), color=FAINT)
arrow(s, Inches(6.85), Inches(3.05), Inches(6.85), Inches(4.1), color=FAINT)
arrow(s, Inches(10.8), Inches(3.05), Inches(10.8), Inches(4.1), color=FAINT)

# Transport legend
add_text(s, Inches(0.9), Inches(5.85), Inches(11.5), Inches(0.35),
         "REALTIME TRANSPORT  —  auto-selecting, in priority order", size=12, color=TEXT, bold=True)
tiers = [
    ("1", "WebSocket → Durable Object", "primary path — the Worker relays every event to all connected peers", BLUE),
    ("2", "Supabase Realtime Broadcast", "fallback when the WS origin isn't configured or reachable", GREEN),
    ("3", "BroadcastChannel", "same-browser, cross-tab rooms with zero network", YELLOW),
]
for i, (n, t, d, c) in enumerate(tiers):
    x = Inches(0.9) + i * Inches(3.95)
    add_rect(s, x, Inches(6.28), Inches(3.65), Inches(0.75), fill=BG_SOFT, line=c, line_w=1.0,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.15)
    step_badge(s, x + Inches(0.15), Inches(6.4), int(n), color=c)
    add_text(s, x + Inches(0.62), Inches(6.3), Inches(2.9), Inches(0.35), t, size=12, color=TEXT, bold=True)
    add_text(s, x + Inches(0.62), Inches(6.62), Inches(2.95), Inches(0.4), d, size=9, color=MUTED, line_spacing=1.05)

# ============================================================================
# SLIDE 9 — FRONTEND ARCHITECTURE
# ============================================================================
s = new_slide()
header(s, "Frontend Architecture", kicker="A modern TanStack application")
front_items = [
    ("File-based routing", "TanStack Router with typed routes — landing, play, lobby, room/:code, leaderboard, profile."),
    ("Server Functions", "getRoom / joinRoom / submitAction are type-safe RPCs the client calls like local functions."),
    ("Authored state, live UI", "Zustand + React Query cache server state; local mirrors for instant, optimistic interaction."),
    ("Component library", "Radix UI primitives, Tailwind design tokens, Framer Motion for card play animations."),
    ("Realtime client", "A transport abstraction sends 'sync' nudges and reacts to peer messages by re-pulling authoritative state."),
]
for i, (t, d) in enumerate(front_items):
    y = Inches(1.6) + i * Inches(1.0)
    add_rect(s, Inches(0.9), y, Inches(11.5), Inches(0.86), fill=BG_CARD, line=BORDER, line_w=0.75,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.09)
    add_rect(s, Inches(0.9), y, Inches(0.06), Inches(0.86), fill=[RED, YELLOW, GREEN, BLUE][i % 4])
    add_text(s, Inches(1.15), y + Inches(0.1), Inches(3.1), Inches(0.7), t, size=13.5, color=TEXT, bold=True, line_spacing=1.05)
    add_text(s, Inches(4.35), y + Inches(0.1), Inches(7.9), Inches(0.7), d, size=11.5, color=MUTED, line_spacing=1.12)

# ============================================================================
# SLIDE 10 — BACKEND ARCHITECTURE
# ============================================================================
s = new_slide()
header(s, "Backend Architecture", kicker="Cloudflare Workers at the edge")
back_items = [
    ("One Durable Object per room", "Rooms are keyed by code and live as globally addressable, "
     "co-located objects — instant, region-aware relay with persistent state."),
    ("WebSocket relay", "The DO accepts sockets, binds identity to each connection, and fans every "
     "message out to the roster. Newcomers and reconnects get a full state snapshot."),
    ("Identity over the wire", "A player can never impersonate another — identity is bound to the "
     "socket, and every action must carry the same playerId."),
    ("Server functions as the API", "Join, start, act and leave are server functions that validate "
     "tokens, enforce rules, persist to D1, and return the authoritative state."),
    ("CORS + upgrade handling", "The Worker handles OPTIONS and WebSocket upgrades with a single, "
     "shared CORS policy."),
]
for i, (t, d) in enumerate(back_items):
    y = Inches(1.6) + i * Inches(1.0)
    add_rect(s, Inches(0.9), y, Inches(11.5), Inches(0.86), fill=BG_CARD, line=BORDER, line_w=0.75,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.09)
    step_badge(s, Inches(1.12), y + Inches(0.26), i + 1, color=[RED, YELLOW, GREEN, BLUE][i % 4])
    add_text(s, Inches(1.7), y + Inches(0.1), Inches(3.6), Inches(0.7), t, size=13.5, color=TEXT, bold=True, line_spacing=1.05)
    add_text(s, Inches(5.4), y + Inches(0.1), Inches(6.85), Inches(0.7), d, size=11.5, color=MUTED, line_spacing=1.1)

# ============================================================================
# SLIDE 11 — REALTIME TRANSPORT LAYER
# ============================================================================
s = new_slide()
header(s, "Realtime Transport Layer", kicker="Reliability through graceful degradation")

# diagram: single client flow, three layers stacked
add_text(s, Inches(0.9), Inches(1.6), Inches(4.6), Inches(0.4), "Every client sends two signals:", size=15, color=TEXT, bold=True)
add_text(s, Inches(0.9), Inches(2.05), Inches(4.6), Inches(1.1),
         [[("→ ", {"color": GREEN}), ("sync nudges", {"bold": True, "color": TEXT}),
           ("  “something changed, re-pull”", {"color": MUTED})],
          [("→ ", {"color": YELLOW}), ("join / leave / hello", {"bold": True, "color": TEXT}),
           ("  roster events", {"color": MUTED})]],
         size=12.5, line_spacing=1.5)

layers = [
    ("LAYER 1", "WebSocket → Durable Object", RED,
     "Primary path. The Worker fans out events to every peer over a persistent socket. "
     "Reconnect-safe: newcomers receive the full latest state."),
    ("LAYER 2", "Supabase Realtime Broadcast", YELLOW,
     "Fallback transport. Broadcasts channeled through Supabase when the WS origin is "
     "not deployed or reachable."),
    ("LAYER 3", "BroadcastChannel (same browser)", GREEN,
     "Local fallback. Cross-tab rooms on the same machine need zero network — perfect "
     "for demos and quick multiplayer."),
]
for i, (l, t, c, d) in enumerate(layers):
    y = Inches(3.05) + i * Inches(1.25)
    add_rect(s, Inches(0.9), y, Inches(11.5), Inches(1.08), fill=BG_CARD, line=c, line_w=1.5,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.09)
    chip(s, Inches(1.15), y + Inches(0.14), Inches(1.05), l, c)
    add_text(s, Inches(2.45), y + Inches(0.12), Inches(4.6), Inches(0.4), t, size=14, color=TEXT, bold=True)
    add_text(s, Inches(1.15), y + Inches(0.52), Inches(11.0), Inches(0.5), d, size=10.5, color=MUTED, line_spacing=1.1)

add_text(s, Inches(0.9), Inches(6.9), Inches(11.5), Inches(0.4),
         "A 3-second poller re-pulls authoritative state regardless of transport, so clients "
         "always converge even if a push is missed.", size=11.5, color=FAINT)

# ============================================================================
# SLIDE 12 — GAME ENGINE
# ============================================================================
s = new_slide()
header(s, "Server-Authoritative Game Engine", kicker="Pure, deterministic TypeScript — zero dependencies")
eng = [
    ("Deterministic core", "The same pure functions create games, apply actions and score rounds — "
     "the server is the single source of truth.", RED),
    ("Strict move validation", "Every play is checked: is it this player's turn, is the card legal, "
     "is a color chosen for wilds, is the +2/+4 stack respected?", YELLOW),
    ("Full UNO ruleset", "Skip, reverse (acts as skip with 2 players), draw-2 and wild+4 stacking "
     "chains, wild color selection, missed-UNO penalty (+2).", GREEN),
    ("Automatic turn timeouts", "A 30s timer auto-draws for idle players so games never stall.", BLUE),
    ("Scoring", "Winners bank the face value of every card left in opponents' hands — numbers, "
     "20 for action cards, 50 for wilds.", RED),
]
for i, (t, d, c) in enumerate(eng):
    y = Inches(1.6) + i * Inches(1.0)
    card(s, Inches(0.9), y, Inches(3.0), Inches(0.86), "", "", accent=[RED, YELLOW, GREEN, BLUE][i % 4])
    add_text(s, Inches(1.12), y + Inches(0.13), Inches(2.6), Inches(0.6), t, size=12.5, color=TEXT, bold=True, line_spacing=1.05)
    add_text(s, Inches(4.15), y + Inches(0.1), Inches(8.1), Inches(0.7), d, size=11.5, color=MUTED, line_spacing=1.1)

# code-ish box
add_rect(s, Inches(0.9), Inches(6.55), Inches(11.5), Inches(0.5), fill=BG_SOFT, line=BORDER, line_w=1.0,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.12)
add_text(s, Inches(1.15), Inches(6.63), Inches(11.0), Inches(0.35),
         [[("engine.ts  ", {"color": GREEN, "font": FONT_MONO}),
           ("createGame() · applyAction() · canPlay() · legalCards() · handleTimeout()",
            {"color": MUTED, "font": FONT_MONO})]],
         size=10.5)

# ============================================================================
# SLIDE 13 — DATA, AUTH & MULTIPLAYER FLOW
# ============================================================================
s = new_slide()
header(s, "Data, Auth & Multiplayer Flow", kicker="Sessions, persistence and live play")

add_text(s, Inches(0.9), Inches(1.55), Inches(5.6), Inches(0.4), "Guest-first identity", size=16, color=TEXT, bold=True)
d_items = [
    ("Guest sessions", "No signup — a local profile (id, name, avatar) is created on first visit."),
    ("JWT per room", "Joining a room returns a token stored in localStorage; every action carries it."),
    ("D1 persistence", "Profiles, game history, wins and leaderboard stats are persisted in SQLite."),
    ("Privacy by design", "Clients only ever receive their own hand — opponents' cards never leave the server."),
]
for i, (t, d) in enumerate(d_items):
    y = Inches(2.05) + i * Inches(1.14)
    add_rect(s, Inches(0.9), y, Inches(5.6), Inches(1.0), fill=BG_CARD, line=BORDER, line_w=0.75,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.09)
    add_text(s, Inches(1.12), y + Inches(0.1), Inches(5.2), Inches(0.35), t, size=12.5, color=TEXT, bold=True)
    add_text(s, Inches(1.12), y + Inches(0.42), Inches(5.2), Inches(0.55), d, size=10.5, color=MUTED, line_spacing=1.1)

# Flow diagram (right)
add_text(s, Inches(7.0), Inches(1.55), Inches(5.5), Inches(0.4), "Live multiplayer loop", size=16, color=TEXT, bold=True)
flow = [
    ("Create / Join", "room code + token", BLUE),
    ("Connect transport", "WS / Supabase / BC", GREEN),
    ("Dispatch action", "play · draw · uno · color", YELLOW),
    ("Engine validates", "turn + legality + rules", RED),
    ("Authoritative state", "fan-out to all peers", RED),
]
fy = Inches(2.05)
for i, (t, d, c) in enumerate(flow):
    add_rect(s, Inches(7.0), fy, Inches(5.0), Inches(0.78), fill=BG_CARD, line=c, line_w=1.25,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.14)
    step_badge(s, Inches(7.18), fy + Inches(0.2), i + 1, color=c)
    add_text(s, Inches(7.62), fy + Inches(0.1), Inches(4.3), Inches(0.35), t, size=12.5, color=TEXT, bold=True)
    add_text(s, Inches(7.62), fy + Inches(0.42), Inches(4.3), Inches(0.35), d, size=9.5, color=MUTED)
    if i < len(flow) - 1:
        arrow(s, Inches(9.5), fy + Inches(0.78), Inches(9.5), fy + Inches(0.86), color=FAINT, w=1.5)
    fy += Inches(0.92)

# ============================================================================
# SLIDE 14 — DEPLOYMENT & CI/CD
# ============================================================================
s = new_slide()
header(s, "Deployment & Continuous Integration", kicker="Ship to the edge, verify every commit")

add_text(s, Inches(0.9), Inches(1.6), Inches(6.0), Inches(0.4), "Deployment targets", size=16, color=TEXT, bold=True)
dep = [
    ("Frontend", "Cloudflare Pages — static build of the React app, edge-cached, globally distributed."),
    ("Backend", "Cloudflare Workers — Durable Object room server deployed via wrangler."),
    ("Realtime", "WebSocket origin on the Worker; Supabase channels as fallback."),
    ("Database", "Cloudflare D1 (SQLite) for profiles, history and leaderboard."),
]
for i, (t, d) in enumerate(dep):
    y = Inches(2.1) + i * Inches(0.98)
    chip(s, Inches(0.9), y + Inches(0.02), Inches(1.75), t.upper(), [RED, YELLOW, GREEN, BLUE][i], filled=True, size=9)
    add_text(s, Inches(2.9), y, Inches(4.6), Inches(0.9), d, size=11, color=MUTED, line_spacing=1.15)

add_text(s, Inches(7.8), Inches(1.6), Inches(5.0), Inches(0.4), "CI pipeline — GitHub Actions", size=16, color=TEXT, bold=True)
ci = [
    ("Typecheck", "tsc --noEmit", "strict, no silent type drift"),
    ("Build", "vite build", "production bundle compiles"),
    ("Lint", "eslint .", "code quality gates"),
]
for i, (t, c, d) in enumerate(ci):
    y = Inches(2.1) + i * Inches(1.25)
    add_rect(s, Inches(7.8), y, Inches(4.6), Inches(1.05), fill=BG_CARD, line=BORDER, line_w=0.75,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.09)
    add_rect(s, Inches(7.8), y, Inches(0.06), Inches(1.05), fill=[BLUE, GREEN, YELLOW][i])
    add_text(s, Inches(8.05), y + Inches(0.12), Inches(2.4), Inches(0.4), t, size=13, color=TEXT, bold=True)
    add_text(s, Inches(10.4), y + Inches(0.12), Inches(1.9), Inches(0.4), c, size=9.5, color=[BLUE, GREEN, YELLOW][i], font=FONT_MONO, align=PP_ALIGN.RIGHT)
    add_text(s, Inches(8.05), y + Inches(0.5), Inches(4.2), Inches(0.4), d, size=10.5, color=MUTED)

card(s, Inches(0.9), Inches(6.05), Inches(11.5), Inches(0.85),
     "One command to deploy", "cd worker && wrangler deploy  —  then point VITE_REALTIME_WS_URL at the "
     "new wss:// origin. The frontend builds and ships independently.", accent=GREEN, title_size=13.5, body_size=11.5)

# ============================================================================
# SLIDE 15 — GAMEPLAY EXPERIENCE (SCREENSHOTS)
# ============================================================================
s = new_slide()
header(s, "Gameplay Experience", kicker="From match setup to the winning card")
p1, h1 = image_centered(s, os.path.join(ASSETS, "play-page.png"), Inches(0.9), Inches(1.9), Inches(4.4))
add_text(s, Inches(0.9), Inches(6.15), Inches(4.4), Inches(0.4),
         "Choose your challenge — bots or a private room", size=11, color=MUTED, align=PP_ALIGN.CENTER)

p2, h2 = image_centered(s, os.path.join(ASSETS, "game-board.png"), Inches(5.85), Inches(1.85), Inches(6.6))
add_text(s, Inches(5.85), Inches(6.15), Inches(6.6), Inches(0.4),
         "Live board — opponents, deck, discard pile, turn timer and move log", size=11, color=MUTED, align=PP_ALIGN.CENTER)

# ============================================================================
# SLIDE 16 — ROOMS & LOBBY (SCREENSHOTS)
# ============================================================================
s = new_slide()
header(s, "Private Rooms & Lobby", kicker="Gather the crew in seconds")
p1, h1 = image_centered(s, os.path.join(ASSETS, "lobby-page.png"), Inches(0.9), Inches(1.85), Inches(5.5))
add_text(s, Inches(0.9), Inches(5.75), Inches(5.5), Inches(0.4),
         "Create or join a room by 6-letter code", size=11, color=MUTED, align=PP_ALIGN.CENTER)

p2, h2 = image_centered(s, os.path.join(ASSETS, "room-lobby.png"), Inches(6.9), Inches(1.85), Inches(5.5))
add_text(s, Inches(6.9), Inches(5.75), Inches(5.5), Inches(0.4),
         "Room lobby — copy invite link, host starts when ready", size=11, color=MUTED, align=PP_ALIGN.CENTER)

add_text(s, Inches(0.9), Inches(6.35), Inches(11.5), Inches(0.6),
         [[("How it works:  ", {"bold": True, "color": TEXT}),
           ("Host opens a room → gets a code + link → shares it → players join from any device → "
            "the host deals the first round.", {"color": MUTED})]],
         size=12, line_spacing=1.2)

# ============================================================================
# SLIDE 17 — STATS & PROGRESSION (SCREENSHOTS)
# ============================================================================
s = new_slide()
header(s, "Stats & Progression", kicker="Every game counts")
p1, h1 = image_centered(s, os.path.join(ASSETS, "leaderboard-page.png"), Inches(0.9), Inches(1.85), Inches(5.5))
add_text(s, Inches(0.9), Inches(5.75), Inches(5.5), Inches(0.4),
         "Global leaderboard — win streaks and totals", size=11, color=MUTED, align=PP_ALIGN.CENTER)

p2, h2 = image_centered(s, os.path.join(ASSETS, "profile-page.png"), Inches(6.9), Inches(1.75), Inches(5.5))
add_text(s, Inches(6.9), Inches(5.85), Inches(5.5), Inches(0.4),
         "Personal profile — games, wins and history", size=11, color=MUTED, align=PP_ALIGN.CENTER)

add_text(s, Inches(0.9), Inches(6.35), Inches(11.5), Inches(0.6),
         [[("No accounts:  ", {"bold": True, "color": TEXT}),
           ("a guest profile follows the player across sessions via local storage; leaderboard stats "
            "are server-persisted in D1.", {"color": MUTED})]],
         size=12, line_spacing=1.2)

# ============================================================================
# SLIDE 18 — PERFORMANCE, SECURITY & RELIABILITY
# ============================================================================
s = new_slide()
header(s, "Performance, Security & Reliability", kicker="Engineering under pressure")
cols = [
    ("Performance", [
        ("Edge-distributed rooms", "Durable Objects run close to players; static assets are served "
         "from the Cloudflare edge cache."),
        ("Slim realtime protocol", "Only 'sync' nudges and roster events travel the wire — state is "
         "pulled once, rendered locally."),
        ("Lazy hydration", "Client routes hydrate only when needed; animations are GPU-friendly."),
    ], GREEN),
    ("Security", [
        ("Server-authoritative", "No client trust: dealing, shuffling, legality and scoring all happen server-side."),
        ("Socket-bound identity", "Actions carry the same playerId as the socket — impersonation is structurally impossible."),
        ("JWT room tokens", "Every room action is authenticated; tokens persist per room in localStorage."),
        ("Hand privacy", "A player never receives opponents' cards — the server only broadcasts results."),
    ], RED),
    ("Reliability", [
        ("Graceful transport fallback", "WS → Supabase → BroadcastChannel with a 3s poller that guarantees convergence."),
        ("Reconnect snapshots", "Late joiners and reconnects receive the full authoritative state instantly."),
        ("CI enforcement", "Typecheck, build and lint gate every commit to prevent regressions."),
    ], BLUE),
]
for i, (title, items, c) in enumerate(cols):
    x = Inches(0.9) + i * Inches(3.95)
    add_text(s, x, Inches(1.6), Inches(3.65), Inches(0.4), title.upper(), size=15, color=c, bold=True)
    yy = Inches(2.08)
    for t, d in items:
        add_rect(s, x, yy, Inches(3.65), Inches(1.08), fill=BG_CARD, line=BORDER, line_w=0.75,
                 shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.08)
        add_rect(s, x, yy, Inches(0.05), Inches(1.08), fill=c)
        add_text(s, x + Inches(0.18), yy + Inches(0.08), Inches(3.3), Inches(0.32), t, size=11, color=TEXT, bold=True)
        add_text(s, x + Inches(0.18), yy + Inches(0.38), Inches(3.35), Inches(0.68), d, size=9, color=MUTED, line_spacing=1.1)
        yy += Inches(1.24)

# ============================================================================
# SLIDE 19 — CHALLENGES & SOLUTIONS
# ============================================================================
s = new_slide()
header(s, "Challenges & Solutions", kicker="Real problems, engineered answers")
chals = [
    ("Realtime that survives anything", "Fragile single-channel WebSockets.",
     "Built an auto-selecting transport (WS → Supabase → BroadcastChannel) plus a 3s state poller, so "
     "clients always converge even when a push is lost."),
    ("Cheat-proof multiplayer", "Client-authoritative games are trivially hackable.",
     "Moved every decision into a deterministic server engine; clients receive only their own hand and "
     "can only nudge peers to re-pull state."),
    ("Concurrency in card games", "Double-plays, stale actions, race conditions.",
     "Room sequence numbers, socket-bound identity, and re-fetch-on-nudge semantics make every action "
     "idempotent and self-healing."),
    ("Type-safety across the wire", "Mismatched client/server contracts.",
     "Zod-validated search params and shared types between client and server functions keep the contract "
     "compile-time checked end to end."),
    ("Keep the build green", "Every commit risks regressions.",
     "CI runs typecheck, build and lint; a strict ESLint config (including the no-control-regex edge "
     "case) keeps the pipeline passing."),
]
for i, (t, prob, sol) in enumerate(chals):
    y = Inches(1.6) + i * Inches(1.04)
    step_badge(s, Inches(0.95), y + Inches(0.02), i + 1, color=[RED, YELLOW, GREEN, BLUE][i % 4])
    add_text(s, Inches(1.5), y, Inches(3.4), Inches(0.4), t, size=13.5, color=TEXT, bold=True)
    add_text(s, Inches(1.5), y + Inches(0.34), Inches(3.3), Inches(0.6),
             [[("Issue: ", {"bold": True, "color": RED}), (prob, {"color": MUTED})]], size=9.5, line_spacing=1.1)
    add_rect(s, Inches(5.05), y - Inches(0.02), Inches(0.03), Inches(0.92), fill=YELLOW)
    add_text(s, Inches(5.3), y + Inches(0.02), Inches(7.1), Inches(0.95),
             [[("Solution: ", {"bold": True, "color": GREEN}), (sol, {"color": MUTED})]], size=10.5, line_spacing=1.18)

# ============================================================================
# SLIDE 20 — FUTURE SCOPE
# ============================================================================
s = new_slide()
header(s, "Future Scope", kicker="Where WildCards goes next")
fut = [
    ("Accounts & OAuth", "Let players claim their guest profile with email or social login and sync stats across devices."),
    ("Matchmaking & ladder", "Ranked queues, ELO-style ratings and seasonal leaderboards for competitive play."),
    ("AI difficulty tiers", "Smarter bots with personality: aggressive, defensive and classic-strategy models."),
    ("Mobile-first & PWA", "Installable offline shell, touch-optimized controls, and push-friendly room invites."),
    ("Tournaments", "Brackets, timed events and custom house rules for organized play."),
    ("Spectator mode", "Watch live rooms with a full move history and commentary."),
]
for i, (t, d) in enumerate(fut):
    x = Inches(0.9) + (i % 3) * Inches(3.95)
    y = Inches(1.6) + (i // 3) * Inches(2.55)
    card(s, x, y, Inches(3.65), Inches(2.3), t, d, accent=[RED, YELLOW, GREEN, BLUE][i % 4], title_size=14, body_size=11.5)

# ============================================================================
# SLIDE 21 — DEMO
# ============================================================================
s = new_slide()
header(s, "Live Demo", kicker="See it in action")
add_text(s, Inches(0.9), Inches(1.75), Inches(11.5), Inches(0.5),
         "What we'll run through:", size=16, color=TEXT, bold=True)
demo = [
    ("Start a game vs bots", "One click — bots fill the table and the first hand is dealt."),
    ("Open a private room", "Create a room, copy the invite link, join from a second tab/device."),
    ("Play a full turn", "Legal plays, wild color choice, draw, and the live turn timer."),
    ("Trigger an edge case", "Say UNO late, start a +2 / +4 chain, and watch the server enforce it."),
    ("Check the scoreboard", "Finish a round and see the profile + leaderboard update."),
]
for i, (t, d) in enumerate(demo):
    y = Inches(2.35) + i * Inches(0.82)
    step_badge(s, Inches(0.95), y, i + 1, color=[RED, YELLOW, GREEN, BLUE][i % 4])
    add_text(s, Inches(1.5), y - Inches(0.03), Inches(4.0), Inches(0.4), t, size=14, color=TEXT, bold=True)
    add_text(s, Inches(5.6), y, Inches(6.8), Inches(0.75), d, size=11.5, color=MUTED, line_spacing=1.15)

card(s, Inches(0.9), Inches(6.35), Inches(11.5), Inches(0.75),
     "Run locally", "pnpm dev  —  then open http://localhost:3000 and pick Play vs Bots or Create a Room.",
     accent=BLUE, title_size=13, body_size=11)

# ============================================================================
# SLIDE 22 — THANK YOU
# ============================================================================
s = new_slide()
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=BG)
q = int(SLIDE_W) // 4
for i, c in enumerate([RED, YELLOW, GREEN, BLUE]):
    add_rect(s, Emu(q * i), 0, Emu(q), Inches(0.12), fill=c)

add_text(s, Inches(0.9), Inches(2.1), Inches(11.5), Inches(0.5),
         "THANK YOU", size=52, color=TEXT, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.9), Inches(3.0), Inches(11.5), Inches(0.5),
         "Questions welcome", size=20, color=YELLOW, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.9), Inches(4.0), Inches(11.5), Inches(0.8),
         "WildCards — browser multiplayer UNO, reimagined for the web.",
         size=14, color=MUTED, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.9), Inches(5.2), Inches(11.5), Inches(0.4),
         "React 19  ·  TypeScript  ·  TanStack Start  ·  Cloudflare Workers  ·  Durable Objects  ·  D1  ·  WebSockets",
         size=12, color=FAINT, align=PP_ALIGN.CENTER)

# ----------------------------------------------------------------------------
prs.save(OUT)
print("Saved:", OUT)
print("Slides:", len(prs.slides._sldIdLst))
